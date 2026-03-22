#!/usr/bin/env python3
"""
Unstructured MCP Server

Converts various document formats to Markdown using Unstructured library.
Supports: PDF, DOCX, PPTX, TXT, MD, HTML, etc.
"""

import asyncio
import json
import os
import sys
import tempfile
import base64
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from unstructured.partition.auto import partition
    from unstructured.chunking.title import chunk_by_title
    from unstructured.documents.elements import Element
    HAS_UNSTRUCTURED = True
except ImportError:
    print("Warning: unstructured not installed. Install with: pip install unstructured", file=sys.stderr)
    HAS_UNSTRUCTURED = False

try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document
    HAS_PYTHON_DOCX = True
except ImportError:
    HAS_PYTHON_DOCX = False

try:
    from pptx import Presentation
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

from mcp.server import Server
from mcp.types import Tool, TextContent

server = Server("unstructured-server")


@server.list_tools()
async def list_tools() -> List[Tool]:
    return [
        Tool(
            name="convert_to_markdown",
            description="Convert a document (PDF, DOCX, PPTX, etc.) to Markdown format",
            inputSchema={
                "type": "object",
                "properties": {
                    "content_base64": {
                        "type": "string",
                        "description": "Base64 encoded file content"
                    },
                    "filename": {
                        "type": "string",
                        "description": "Original filename (used to detect format)"
                    },
                    "chunking_strategy": {
                        "type": "string",
                        "enum": ["none", "by_title", "basic"],
                        "description": "Chunking strategy for long documents",
                        "default": "by_title"
                    },
                    "max_characters": {
                        "type": "integer",
                        "description": "Max characters per chunk",
                        "default": 1500
                    },
                    "extract_images": {
                        "type": "boolean",
                        "description": "Whether to extract and describe images",
                        "default": False
                    }
                },
                "required": ["content_base64", "filename"]
            }
        ),
        Tool(
            name="batch_convert",
            description="Batch convert multiple documents to Markdown",
            inputSchema={
                "type": "object",
                "properties": {
                    "files": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "content_base64": {"type": "string"},
                                "filename": {"type": "string"}
                            },
                            "required": ["content_base64", "filename"]
                        },
                        "description": "List of files to convert"
                    },
                    "chunking_strategy": {
                        "type": "string",
                        "default": "by_title"
                    },
                    "max_characters": {
                        "type": "integer",
                        "default": 1500
                    }
                },
                "required": ["files"]
            }
        ),
        Tool(
            name="get_supported_formats",
            description="Get list of supported document formats",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        )
    ]


@server.call_tool()
async def call_tool(name: str, arguments: Dict[str, Any]) -> List[TextContent]:
    try:
        if name == "convert_to_markdown":
            result = await convert_single_document(arguments)
            return [TextContent(type="text", text=json.dumps(result, indent=2))]
        
        elif name == "batch_convert":
            result = await batch_convert_documents(arguments)
            return [TextContent(type="text", text=json.dumps(result, indent=2))]
        
        elif name == "get_supported_formats":
            result = get_supported_formats()
            return [TextContent(type="text", text=json.dumps(result, indent=2))]
        
        else:
            return [TextContent(type="text", text=json.dumps({"error": f"Unknown tool: {name}"}))]
    
    except Exception as e:
        return [TextContent(type="text", text=json.dumps({
            "error": str(e),
            "type": type(e).__name__
        }))]


async def convert_single_document(args: Dict[str, Any]) -> Dict[str, Any]:
    """Convert a single document to Markdown."""
    content_base64 = args["content_base64"]
    filename = args["filename"]
    chunking_strategy = args.get("chunking_strategy", "by_title")
    max_characters = args.get("max_characters", 1500)
    extract_images = args.get("extract_images", False)
    
    # Decode content
    try:
        content = base64.b64decode(content_base64)
    except Exception as e:
        return {"error": f"Failed to decode base64 content: {str(e)}"}
    
    # Detect file extension
    ext = Path(filename).suffix.lower()
    
    # Choose conversion method based on extension and available libraries
    if ext == ".pdf":
        markdown, metadata = await convert_pdf(content, extract_images)
    elif ext in [".docx", ".doc"]:
        markdown, metadata = await convert_docx(content)
    elif ext in [".pptx", ".ppt"]:
        markdown, metadata = await convert_pptx(content)
    elif ext in [".md", ".txt", ".markdown"]:
        markdown, metadata = convert_text(content, filename)
    else:
        # Try unstructured for other formats
        markdown, metadata = await convert_with_unstructured(
            content, filename, chunking_strategy, max_characters
        )
    
    # Apply chunking if requested
    chunks = None
    if chunking_strategy != "none" and HAS_UNSTRUCTURED:
        chunks = chunk_markdown(markdown, max_characters)
    
    return {
        "filename": filename,
        "markdown": markdown,
        "chunks": chunks,
        "chunk_count": len(chunks) if chunks else None,
        "metadata": metadata,
        "word_count": len(markdown.split()),
        "char_count": len(markdown)
    }


async def convert_pdf(content: bytes, extract_images: bool = False) -> tuple:
    """Convert PDF to Markdown using PyMuPDF."""
    if not HAS_PYMUPDF:
        # Fall back to unstructured
        return await convert_with_unstructured(content, "document.pdf", "by_title", 1500)
    
    markdown_parts = []
    metadata = {"pages": 0, "images": 0}
    
    # Create temporary file
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name
    
    try:
        doc = fitz.open(tmp_path)
        metadata["pages"] = doc.page_count
        
        for page_num in range(doc.page_count):
            page = doc[page_num]
            text = page.get_text("text")
            
            if text.strip():
                markdown_parts.append(f"## Page {page_num + 1}\n\n{text}\n")
            
            if extract_images:
                images = page.get_images()
                metadata["images"] += len(images)
                # Could add image extraction logic here
        
        doc.close()
    finally:
        os.unlink(tmp_path)
    
    return "\n".join(markdown_parts), metadata


async def convert_docx(content: bytes) -> tuple:
    """Convert DOCX to Markdown."""
    if not HAS_PYTHON_DOCX:
        return await convert_with_unstructured(content, "document.docx", "by_title", 1500)
    
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name
    
    try:
        doc = Document(tmp_path)
        markdown_parts = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Check if heading
                if para.style.name.startswith('Heading'):
                    level = para.style.name.split()[-1]
                    try:
                        level_num = int(level)
                        markdown_parts.append(f"{'#' * level_num} {text}\n")
                    except:
                        markdown_parts.append(f"{text}\n")
                else:
                    markdown_parts.append(f"{text}\n")
        
        metadata = {"paragraphs": len(doc.paragraphs)}
        return "\n".join(markdown_parts), metadata
    
    finally:
        os.unlink(tmp_path)


async def convert_pptx(content: bytes) -> tuple:
    """Convert PPTX to Markdown."""
    if not HAS_PYTHON_PPTX:
        return await convert_with_unstructured(content, "document.pptx", "by_title", 1500)
    
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name
    
    try:
        prs = Presentation(tmp_path)
        markdown_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_parts.append(f"## Slide {slide_num}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    markdown_parts.append(shape.text + "\n")
            
            markdown_parts.append("\n")
        
        metadata = {"slides": len(prs.slides)}
        return "\n".join(markdown_parts), metadata
    
    finally:
        os.unlink(tmp_path)


def convert_text(content: bytes, filename: str) -> tuple:
    """Convert plain text to Markdown (passthrough)."""
    text = content.decode('utf-8', errors='ignore')
    metadata = {
        "lines": len(text.split('\n')),
        "original_format": Path(filename).suffix
    }
    return text, metadata


async def convert_with_unstructured(
    content: bytes, 
    filename: str, 
    chunking_strategy: str,
    max_characters: int
) -> tuple:
    """Convert using Unstructured library."""
    if not HAS_UNSTRUCTURED:
        # Last resort: treat as text
        return convert_text(content, filename)
    
    with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name
    
    try:
        elements = partition(filename=tmp_path)
        
        if chunking_strategy == "by_title":
            chunks = chunk_by_title(elements, max_characters=max_characters)
            markdown = "\n\n".join([str(chunk) for chunk in chunks])
        else:
            markdown = "\n\n".join([str(elem) for elem in elements])
        
        metadata = {
            "elements": len(elements),
            "method": "unstructured"
        }
        
        return markdown, metadata
    
    finally:
        os.unlink(tmp_path)


def chunk_markdown(markdown: str, max_chars: int) -> List[str]:
    """Split markdown into chunks."""
    if len(markdown) <= max_chars:
        return [markdown]
    
    chunks = []
    lines = markdown.split('\n')
    current_chunk = []
    current_size = 0
    
    for line in lines:
        line_size = len(line) + 1  # +1 for newline
        
        if current_size + line_size > max_chars and current_chunk:
            chunks.append('\n'.join(current_chunk))
            current_chunk = []
            current_size = 0
        
        current_chunk.append(line)
        current_size += line_size
    
    if current_chunk:
        chunks.append('\n'.join(current_chunk))
    
    return chunks


async def batch_convert_documents(args: Dict[str, Any]) -> Dict[str, Any]:
    """Batch convert multiple documents."""
    files = args["files"]
    chunking_strategy = args.get("chunking_strategy", "by_title")
    max_characters = args.get("max_characters", 1500)
    
    results = []
    success_count = 0
    fail_count = 0
    
    for file_info in files:
        try:
            result = await convert_single_document({
                "content_base64": file_info["content_base64"],
                "filename": file_info["filename"],
                "chunking_strategy": chunking_strategy,
                "max_characters": max_characters
            })
            
            if "error" not in result:
                success_count += 1
            else:
                fail_count += 1
            
            results.append(result)
        except Exception as e:
            fail_count += 1
            results.append({
                "filename": file_info.get("filename", "unknown"),
                "error": str(e)
            })
    
    return {
        "total": len(files),
        "success": success_count,
        "failed": fail_count,
        "results": results
    }


def get_supported_formats() -> Dict[str, Any]:
    """Get list of supported formats."""
    formats = {
        "primary": [
            {"extension": ".pdf", "library": "PyMuPDF (fitz)", "status": "available" if HAS_PYMUPDF else "unavailable"},
            {"extension": ".docx", "library": "python-docx", "status": "available" if HAS_PYTHON_DOCX else "unavailable"},
            {"extension": ".pptx", "library": "python-pptx", "status": "available" if HAS_PYTHON_PPTX else "unavailable"},
            {"extension": ".txt", "library": "built-in", "status": "available"},
            {"extension": ".md", "library": "built-in", "status": "available"},
        ],
        "unstructured": [
            ".html", ".htm", ".xml", ".json", ".csv",
            ".eml", ".msg", ".rst", ".rtf", ".tsv"
        ],
        "installation": {
            "pymupdf": "pip install PyMuPDF",
            "python-docx": "pip install python-docx",
            "python-pptx": "pip install python-pptx",
            "unstructured": "pip install unstructured[pdf,docx,pptx]"
        }
    }
    
    return formats


async def main():
    from mcp.server.stdio import stdio_server
    
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            server.create_initialization_options()
        )


if __name__ == "__main__":
    asyncio.run(main())